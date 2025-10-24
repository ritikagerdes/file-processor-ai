"""
File processing module for handling various file types and creating embeddings.
"""

import hashlib
import io
import json
import logging
import magic
import os
import uuid
from typing import Any, Dict, List, Optional, Tuple, Union

import openai
import PyPDF2
import tiktoken
from bs4 import BeautifulSoup
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from app.config import settings
from app.core.aws_client import aws_manager
from app.core.database import db_manager, File, FileChunk
from app.core.security import security_manager

logger = logging.getLogger(__name__)


class FileProcessingError(Exception):
    """File processing error."""
    pass


class UnsupportedFileTypeError(FileProcessingError):
    """Unsupported file type error."""
    pass


class FileProcessor:
    """Handles file processing, chunking, and embedding generation."""
    
    def __init__(self):
        """Initialize file processor with OpenAI client and text splitter."""
        self.openai_client = openai.OpenAI(api_key=settings.openai_api_key)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            length_function=len,
        )
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
    
    def detect_file_type(self, file_data: bytes, filename: str) -> str:
        """
        Detect file type using magic library and filename extension.
        
        Args:
            file_data: File data as bytes
            filename: Original filename
            
        Returns:
            Detected file type
            
        Raises:
            UnsupportedFileTypeError: If file type is not supported
        """
        try:
            # Use python-magic to detect file type
            mime_type = magic.from_buffer(file_data, mime=True)
            
            # Map MIME types to our supported types
            mime_to_type = {
                'text/plain': 'txt',
                'text/html': 'html',
                'text/markdown': 'md',
                'application/pdf': 'pdf',
                'application/msword': 'doc',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
                'application/vnd.ms-excel': 'xls',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
                'application/vnd.ms-powerpoint': 'ppt',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
                'application/json': 'json',
                'application/xml': 'xml',
                'text/csv': 'csv',
            }
            
            file_type = mime_to_type.get(mime_type)
            if not file_type:
                # Fallback to file extension
                extension = filename.split('.')[-1].lower()
                supported_extensions = {
                    'txt', 'md', 'html', 'htm', 'pdf', 'doc', 'docx', 
                    'xls', 'xlsx', 'ppt', 'pptx', 'json', 'xml', 'csv'
                }
                if extension in supported_extensions:
                    file_type = extension
                else:
                    raise UnsupportedFileTypeError(f"Unsupported file type: {mime_type} or extension: {extension}")
            
            return file_type
        except Exception as e:
            logger.error(f"Failed to detect file type: {e}")
            raise UnsupportedFileTypeError(f"Failed to detect file type: {str(e)}") from e
    
    def extract_text_from_file(self, file_data: bytes, file_type: str, filename: str) -> str:
        """
        Extract text content from various file types.
        
        Args:
            file_data: File data as bytes
            file_type: Detected file type
            filename: Original filename
            
        Returns:
            Extracted text content
            
        Raises:
            FileProcessingError: If text extraction fails
        """
        try:
            if file_type in ['txt', 'md', 'html', 'htm', 'json', 'xml', 'csv']:
                return self._extract_text_from_text_file(file_data, file_type)
            elif file_type == 'pdf':
                return self._extract_text_from_pdf(file_data)
            elif file_type in ['doc', 'docx']:
                return self._extract_text_from_word(file_data, file_type)
            elif file_type in ['xls', 'xlsx']:
                return self._extract_text_from_excel(file_data, file_type)
            elif file_type in ['ppt', 'pptx']:
                return self._extract_text_from_powerpoint(file_data, file_type)
            else:
                raise UnsupportedFileTypeError(f"Text extraction not supported for file type: {file_type}")
        except Exception as e:
            logger.error(f"Failed to extract text from {file_type} file: {e}")
            raise FileProcessingError(f"Text extraction failed: {str(e)}") from e
    
    def _extract_text_from_text_file(self, file_data: bytes, file_type: str) -> str:
        """Extract text from plain text files."""
        try:
            text = file_data.decode('utf-8')
            
            if file_type in ['html', 'htm']:
                soup = BeautifulSoup(text, 'html.parser')
                text = soup.get_text()
            elif file_type == 'json':
                # Pretty print JSON for better readability
                json_data = json.loads(text)
                text = json.dumps(json_data, indent=2)
            elif file_type == 'xml':
                soup = BeautifulSoup(text, 'xml')
                text = soup.get_text()
            
            return text
        except UnicodeDecodeError:
            # Try with different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    return file_data.decode(encoding)
                except UnicodeDecodeError:
                    continue
            raise FileProcessingError("Unable to decode file with any supported encoding")
    
    def _extract_text_from_pdf(self, file_data: bytes) -> str:
        """Extract text from PDF files."""
        try:
            pdf_file = io.BytesIO(file_data)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text
        except Exception as e:
            raise FileProcessingError(f"PDF text extraction failed: {str(e)}") from e
    
    def _extract_text_from_word(self, file_data: bytes, file_type: str) -> str:
        """Extract text from Word documents."""
        try:
            if file_type == 'docx':
                doc_file = io.BytesIO(file_data)
                doc = Document(doc_file)
                text = ""
                
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                
                return text
            else:
                # For .doc files, we would need python-docx2txt or similar
                raise FileProcessingError("DOC file support requires additional dependencies")
        except Exception as e:
            raise FileProcessingError(f"Word document text extraction failed: {str(e)}") from e
    
    def _extract_text_from_excel(self, file_data: bytes, file_type: str) -> str:
        """Extract text from Excel files."""
        try:
            excel_file = io.BytesIO(file_data)
            workbook = load_workbook(excel_file, data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                text += "\n"
            
            return text
        except Exception as e:
            raise FileProcessingError(f"Excel text extraction failed: {str(e)}") from e
    
    def _extract_text_from_powerpoint(self, file_data: bytes, file_type: str) -> str:
        """Extract text from PowerPoint files."""
        try:
            if file_type == 'pptx':
                ppt_file = io.BytesIO(file_data)
                presentation = Presentation(ppt_file)
                text = ""
                
                for slide_num, slide in enumerate(presentation.slides, 1):
                    text += f"Slide {slide_num}:\n"
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    text += "\n"
                
                return text
            else:
                raise FileProcessingError("PPT file support requires additional dependencies")
        except Exception as e:
            raise FileProcessingError(f"PowerPoint text extraction failed: {str(e)}") from e
    
    def chunk_text(self, text: str) -> List[str]:
        """
        Split text into chunks for processing.
        
        Args:
            text: Text to chunk
            
        Returns:
            List of text chunks
        """
        try:
            chunks = self.text_splitter.split_text(text)
            return [chunk.strip() for chunk in chunks if chunk.strip()]
        except Exception as e:
            logger.error(f"Failed to chunk text: {e}")
            raise FileProcessingError(f"Text chunking failed: {str(e)}") from e
    
    def generate_embedding(self, text: str) -> List[float]:
        """
        Generate embedding for text using OpenAI API.
        
        Args:
            text: Text to embed
            
        Returns:
            Embedding vector
            
        Raises:
            FileProcessingError: If embedding generation fails
        """
        try:
            # Use OpenAI's embedding API
            response = self.openai_client.embeddings.create(
                model=settings.embedding_model,
                input=text
            )
            return response.data[0].embedding
        except Exception as e:
            logger.error(f"Failed to generate embedding: {e}")
            raise FileProcessingError(f"Embedding generation failed: {str(e)}") from e
    
    def process_file(self, file_data: bytes, filename: str, client_id: str) -> Dict[str, Any]:
        """
        Process a file: extract text, chunk, generate embeddings, and store.
        
        Args:
            file_data: File data as bytes
            filename: Original filename
            client_id: Client identifier
            
        Returns:
            Processing result dictionary
            
        Raises:
            FileProcessingError: If processing fails
        """
        try:
            # Detect file type
            file_type = self.detect_file_type(file_data, filename)
            
            # Extract text
            text_content = self.extract_text_from_file(file_data, file_type, filename)
            
            if not text_content.strip():
                raise FileProcessingError("No text content found in file")
            
            # Generate secure filename
            secure_filename = security_manager.generate_secure_filename(filename, client_id)
            
            # Create file hash for integrity verification
            content_hash = security_manager.create_audit_hash(text_content)
            
            # Upload file to S3
            s3_key = f"clients/{client_id}/files/{secure_filename}"
            s3_url = aws_manager.get_s3_client().upload_file(
                file_data, 
                s3_key,
                metadata={
                    'original_filename': filename,
                    'client_id': client_id,
                    'file_type': file_type,
                    'content_hash': content_hash
                }
            )
            
            # Create file record in database
            file_record = db_manager.create_file({
                'id': str(uuid.uuid4()),
                'client_id': client_id,
                'original_filename': filename,
                'secure_filename': secure_filename,
                'file_type': file_type,
                'file_size': len(file_data),
                's3_key': s3_key,
                'content_hash': content_hash,
                'is_processed': False
            })
            
            # Chunk the text
            chunks = self.chunk_text(text_content)
            
            # Process chunks and create embeddings
            chunk_records = []
            for i, chunk_text in enumerate(chunks):
                try:
                    # Generate embedding
                    embedding = self.generate_embedding(chunk_text)
                    
                    # Encrypt embedding for storage
                    encrypted_embedding = security_manager.encrypt_data(json.dumps(embedding))
                    
                    # Create chunk hash
                    chunk_hash = security_manager.create_audit_hash(chunk_text)
                    
                    chunk_record = {
                        'id': str(uuid.uuid4()),
                        'file_id': file_record.id,
                        'chunk_index': i,
                        'chunk_text': chunk_text,
                        'chunk_embedding': encrypted_embedding,
                        'chunk_hash': chunk_hash
                    }
                    chunk_records.append(chunk_record)
                    
                except Exception as e:
                    logger.warning(f"Failed to process chunk {i}: {e}")
                    continue
            
            # Save chunks to database
            if chunk_records:
                db_manager.create_file_chunks(chunk_records)
                db_manager.update_file_processing_status(file_record.id, True)
            
            return {
                'file_id': file_record.id,
                'filename': filename,
                'file_type': file_type,
                'file_size': len(file_data),
                'chunks_created': len(chunk_records),
                's3_url': s3_url,
                'status': 'processed'
            }
            
        except Exception as e:
            logger.error(f"File processing failed: {e}")
            raise FileProcessingError(f"File processing failed: {str(e)}") from e
    
    def get_file_chunks_for_search(self, file_id: str) -> List[Dict[str, Any]]:
        """
        Get file chunks for search operations.
        
        Args:
            file_id: File identifier
            
        Returns:
            List of chunk data with decrypted embeddings
        """
        try:
            chunks = db_manager.get_file_chunks(file_id)
            chunk_data = []
            
            for chunk in chunks:
                # Decrypt embedding
                decrypted_embedding = security_manager.decrypt_data(chunk.chunk_embedding)
                embedding = json.loads(decrypted_embedding)
                
                chunk_data.append({
                    'id': chunk.id,
                    'chunk_index': chunk.chunk_index,
                    'chunk_text': chunk.chunk_text,
                    'embedding': embedding,
                    'chunk_hash': chunk.chunk_hash
                })
            
            return chunk_data
        except Exception as e:
            logger.error(f"Failed to get file chunks: {e}")
            raise FileProcessingError(f"Failed to get file chunks: {str(e)}") from e


# Global file processor instance
file_processor = FileProcessor()